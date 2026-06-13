"""文档加载器——策略模式 + 组合模式。

每种文档类型对应一个策略类，FallbackLoaderStrategy 组合多个策略形成优先级链。
"""

import logging
import os
import shutil
from abc import ABC, abstractmethod

from langchain_core.documents import Document

logger = logging.getLogger(__name__)


class DocumentLoaderStrategy(ABC):
    """文档加载策略抽象接口。"""

    @abstractmethod
    def load(self, file_path: str) -> list[Document]:
        """加载文档并返回 LangChain Document 列表。"""
        ...


class PyPDFLoaderStrategy(DocumentLoaderStrategy):
    """PDF 加载——langchain_community PyPDFLoader（备选策略）。"""

    def load(self, file_path: str) -> list[Document]:
        from langchain_community.document_loaders import PyPDFLoader
        return PyPDFLoader(file_path).load()


class OpenDataLoaderPDFStrategy(DocumentLoaderStrategy):
    """PDF 加载——langchain-opendataloader-pdf（优先策略）。"""

    def load(self, file_path: str) -> list[Document]:
        if not os.path.isfile(file_path):
            raise FileNotFoundError(f"PDF file not found: {file_path}")
        from langchain_opendataloader_pdf import OpenDataLoaderPDFLoader
        docs = OpenDataLoaderPDFLoader(
                    file_path=str(file_path),
                    format="markdown",
                    table_method="cluster",
                    include_header_footer=True,
                    image_output="external",
                    image_dir="./images",
                    image_format="png"
                ).load()
        if not docs:
            raise RuntimeError(
                f"OpenDataLoaderPDF returned empty result for {file_path}"
            )
        return docs


class DocxLoaderStrategy(DocumentLoaderStrategy):
    """Word 文档加载。"""

    def load(self, file_path: str) -> list[Document]:
        from langchain_community.document_loaders import Docx2txtLoader
        return Docx2txtLoader(file_path).load()


class UnstructuredExcelStrategy(DocumentLoaderStrategy):
    """Excel 表格加载。"""

    def load(self, file_path: str) -> list[Document]:
        from langchain_community.document_loaders import UnstructuredExcelLoader
        return UnstructuredExcelLoader(file_path).load()


class PythonPPTXLoaderStrategy(DocumentLoaderStrategy):
    """PowerPoint 加载——使用 python-pptx（跨平台，无需 unstructured）。"""

    def load(self, file_path: str) -> list[Document]:
        from pptx import Presentation

        prs = Presentation(file_path)
        docs: list[Document] = []
        for i, slide in enumerate(prs.slides):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                docs.append(Document(
                    page_content="\n".join(texts),
                    metadata={"slide": i, "source": str(file_path)},
                ))
        if not docs:
            raise RuntimeError(f"No text content extracted from {file_path}")
        return docs


class UnstructuredPPTStrategy(DocumentLoaderStrategy):
    """PowerPoint 加载——unstructured（非 Windows 备选，Windows 上会崩溃）。"""

    def load(self, file_path: str) -> list[Document]:
        import platform
        if platform.system() == "Windows":
            raise RuntimeError(
                "UnstructuredPowerPointLoader is not supported on Windows "
                "due to python-magic incompatibility"
            )
        from langchain_community.document_loaders import UnstructuredPowerPointLoader
        return UnstructuredPowerPointLoader(file_path).load()


class CSVLoaderStrategy(DocumentLoaderStrategy):
    """CSV 加载。"""

    def load(self, file_path: str) -> list[Document]:
        from langchain_community.document_loaders import CSVLoader
        return CSVLoader(file_path).load()


class JSONLoaderStrategy(DocumentLoaderStrategy):
    """JSON 加载。"""

    def load(self, file_path: str) -> list[Document]:
        from langchain_community.document_loaders import JSONLoader
        return JSONLoader(file_path, jq_schema=".", text_content=False).load()


class MarkdownLoaderStrategy(DocumentLoaderStrategy):
    """Markdown 加载。"""

    def load(self, file_path: str) -> list[Document]:
        from langchain_community.document_loaders import UnstructuredMarkdownLoader
        return UnstructuredMarkdownLoader(file_path, mode="elements").load()


class HTMLLoaderStrategy(DocumentLoaderStrategy):
    """HTML 加载——使用 BeautifulSoup 解析并提取纯文本。"""

    def load(self, file_path: str) -> list[Document]:
        from langchain_community.document_loaders import BSHTMLLoader
        return BSHTMLLoader(file_path, open_encoding="utf-8").load()


class TextLoaderStrategy(DocumentLoaderStrategy):
    """纯文本加载。"""

    def load(self, file_path: str) -> list[Document]:
        from langchain_community.document_loaders import TextLoader
        return TextLoader(file_path, encoding="utf-8").load()


class ImageLoaderStrategy(DocumentLoaderStrategy):
    """图片加载（OCR/多模态）。需要 Tesseract OCR 已安装并在 PATH 中。"""

    def load(self, file_path: str) -> list[Document]:
        if not shutil.which("tesseract"):
            raise RuntimeError(
                "Tesseract OCR is not installed or not in PATH. "
                "Image loading requires Tesseract for text extraction."
            )
        from langchain_community.document_loaders import UnstructuredImageLoader
        return UnstructuredImageLoader(file_path).load()


class FallbackLoaderStrategy(DocumentLoaderStrategy):
    """组合策略——按优先级链依次尝试，返回第一个成功的结果。"""

    def __init__(self, strategies: list[DocumentLoaderStrategy]):
        self._strategies = strategies

    def load(self, file_path: str) -> list[Document]:
        errors: list[str] = []
        for strategy in self._strategies:
            try:
                return strategy.load(file_path)
            except Exception as e:
                errors.append(f"{type(strategy).__name__}: {e}")
                logger.debug("Loader %s failed: %s", type(strategy).__name__, e)
        raise ValueError(f"All loaders failed: {'; '.join(errors)}")


class DocumentLoaderRegistry:
    """文档加载策略注册表。"""

    def __init__(self):
        self._strategies: dict[str, DocumentLoaderStrategy] = {}

    def register(self, file_type: str, strategy: DocumentLoaderStrategy) -> None:
        self._strategies[file_type] = strategy

    def get_strategy(self, file_type: str) -> DocumentLoaderStrategy:
        if file_type not in self._strategies:
            raise ValueError(f"Unsupported file type: {file_type}")
        return self._strategies[file_type]

    def load(self, file_path: str, file_type: str) -> list[Document]:
        return self.get_strategy(file_type).load(file_path)


def create_default_loader_registry() -> DocumentLoaderRegistry:
    """创建预注册所有内置文件类型的加载器注册表。"""
    registry = DocumentLoaderRegistry()

    # PDF: 优先 opendataloader_pdf，备选 PyPDFLoader
    registry.register("pdf", FallbackLoaderStrategy([
        OpenDataLoaderPDFStrategy(),
        PyPDFLoaderStrategy(),
    ]))

    registry.register("docx", DocxLoaderStrategy())
    registry.register("xlsx", UnstructuredExcelStrategy())
    registry.register("pptx", FallbackLoaderStrategy([
        PythonPPTXLoaderStrategy(),
        UnstructuredPPTStrategy(),
    ]))
    registry.register("csv", CSVLoaderStrategy())
    registry.register("json", JSONLoaderStrategy())
    registry.register("md", MarkdownLoaderStrategy())
    registry.register("html", HTMLLoaderStrategy())
    registry.register("txt", TextLoaderStrategy())

    image_strategy = ImageLoaderStrategy()
    for ft in ("png", "jpg", "jpeg", "image"):
        registry.register(ft, image_strategy)

    return registry
